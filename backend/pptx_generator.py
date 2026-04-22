# PPTX Generator module built with python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def generate_claim_chart_pptx(claims_data: list, output_filename: str):
    """
    Generates the PPTX file.
    Each independent claim gets one slide.
    Left side: Claim text with numbered elements.
    Right side: Selected figure image.
    claims_data = [
        {
            "claim_text": "...",
            "best_figure_id": "Page 1 Figure 0",
            "image_path": "path/to/img",
            "elements": [{"text": "...", "numeral": "10"}]
        }, ...
    ]
    """
    
    prs = Presentation()
    
    # 6 is the layout for Title and Two Content or just blank
    blank_slide_layout = prs.slide_layouts[6]
    
    for claim in claims_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Left Panel (Text box)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(4.5)
        height = Inches(6.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Independent Claim Elements"
        p.font.bold = True
        p.font.size = Pt(18)
        
        for el in claim.get("elements", []):
            p = tf.add_paragraph()
            p.text = f"• {el.get('text', '')} ({el.get('numeral', '')})"
            p.font.size = Pt(14)
            
        p = tf.add_paragraph()
        p.text = f"\nFull Claim:\n{claim.get('claim_text', '')}"
        p.font.size = Pt(12)
        
        # Right Panel (Image)
        img_path = claim.get("image_path")
        if img_path:
            left_img = Inches(5.2)
            top_img = Inches(1.0)
            height_img = Inches(5.5)
            try:
                slide.shapes.add_picture(img_path, left_img, top_img, height=height_img)
            except Exception as e:
                print(f"Error adding image {img_path} to slide: {e}")
        else:
             txBox2 = slide.shapes.add_textbox(Inches(5.2), Inches(3.0), Inches(4.0), Inches(1.0))
             txBox2.text_frame.text = "No matching figure found"
        
    prs.save(output_filename)
    return output_filename
